"""Document processing module for loading and chunking documents."""

from typing import List
import os
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain_core.documents import Document


class DocumentProcessor:
    """Handles document loading and chunking operations."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: Size of each text chunk
            chunk_overlap: Overlap between consecutive chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    def _load_pptx(self, file_path: str) -> List[Document]:
        """
        Load PPTX file using python-pptx directly to avoid NLTK dependencies.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            List of Document objects
        """
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide in prs.slides:
                slide_text = []
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                # Join slide content
                if slide_text:
                    text_content.append("\n".join(slide_text))
            
            # Join all slides with double newlines
            full_text = "\n\n".join(text_content)
            
            # Create a single document for the whole presentation
            # The chunker will split it later
            metadata = {"source": file_path}
            return [Document(page_content=full_text, metadata=metadata)]
            
        except Exception as e:
            raise ValueError(f"Failed to load PPTX file: {e}")

    def load_document(self, file_path: str) -> List[Document]:
        """
        Load a document from file path.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of Document objects
            
        Raises:
            ValueError: If file type is not supported
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            loader = PyPDFLoader(file_path)
            return loader.load()
        elif file_extension == '.txt':
            loader = TextLoader(file_path)
            return loader.load()
        elif file_extension == '.docx':
            loader = Docx2txtLoader(file_path)
            return loader.load()
        elif file_extension == '.pptx':
            return self._load_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    def chunk_documents(self, documents: List[Document]) -> List[Document]:
        """
        Split documents into smaller chunks.
        
        Args:
            documents: List of Document objects to chunk
            
        Returns:
            List of chunked Document objects
        """
        return self.text_splitter.split_documents(documents)
    
    def process_document(self, file_path: str) -> List[Document]:
        """
        Complete processing pipeline: load and chunk document.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of chunked Document objects
        """
        documents = self.load_document(file_path)
        chunks = self.chunk_documents(documents)
        return chunks
